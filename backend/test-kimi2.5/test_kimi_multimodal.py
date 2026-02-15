import os
import time
import sys
import csv
import json
import base64
import requests
from openai import OpenAI
try:
    import dashscope
    from dashscope.audio.asr import Recognition, RecognitionCallback, RecognitionResult
except ImportError:
    dashscope = None
    # 定义虚拟类以防 ImportError 时后续代码报错
    class RecognitionCallback: pass
    class RecognitionResult: pass

# 尝试导入生成文件的库，如果不存在则在生成时跳过
try:
    from fpdf import FPDF
except ImportError:
    FPDF = None

try:
    from docx import Document
except ImportError:
    Document = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

# -----------------------------------------------------------------------------
# 配置区域
# -----------------------------------------------------------------------------
# 请设置环境变量 MOONSHOT_MOONSHOT_API_KEY，或者直接在此处填入您的 API Key
MOONSHOT_API_KEY = "sk-a62knARhMuUNiGsBSEnumZ2BZt2AVo8NzeI8CbybsELV3Rs9"
MOONSHOT_BASE_URL = "https://api.moonshot.cn/v1"

DASHSCOPE_API_KEY = "sk-c4383b80d4ce4bd1a20feefa5b9e9a6f"


# 用户指定的模型名称
MODEL_NAME = "kimi-k2.5"

# -----------------------------------------------------------------------------
# 初始化客户端
# -----------------------------------------------------------------------------
if DASHSCOPE_API_KEY == "":
    print("请先设置环境变量 DASHSCOPE_API_KEY 或在代码中填入您的 API Key。")
    # 为了演示方便，我们允许脚本继续运行，但 API 调用会失败
    pass

if not MOONSHOT_API_KEY:
        print("请先设置环境变量 MOONSHOT_API_KEY 或在代码中填入您的 API Key。")
        # 为了演示方便，我们允许脚本继续运行，但 API 调用会失败
        pass

try:
    client = OpenAI(
        api_key=MOONSHOT_API_KEY,
        base_url=MOONSHOT_BASE_URL,
    )
except Exception as e:
    print(f"初始化 OpenAI 客户端失败: {e}")
    sys.exit(1)

# -----------------------------------------------------------------------------
# 功能函数
# -----------------------------------------------------------------------------

def encode_image_to_base64(image_source):
    """
    将图片转换为 Base64 编码字符串
    :param image_source: 图片路径或 URL
    :return: base64 string
    """
    try:
        # 如果是 URL
        if image_source.startswith("http"):
            response = requests.get(image_source)
            if response.status_code == 200:
                return base64.b64encode(response.content).decode('utf-8')
            else:
                print(f"❌ 下载图片失败: {response.status_code}")
                return None
        # 如果是本地文件
        elif os.path.exists(image_source):
            with open(image_source, "rb") as image_file:
                return base64.b64encode(image_file.read()).decode('utf-8')
        else:
            print(f"❌ 图片文件不存在: {image_source}")
            return None
    except Exception as e:
        print(f"❌ 图片编码错误: {e}")
        return None

def test_image_understanding():
    """
    测试图片理解功能 (Vision)
    使用 Kimi 模型处理图片 (Base64)
    """
    print("\n" + "="*50)
    print("开始测试：图片理解 (Image Understanding)")
    print("="*50)

    # 示例图片 URL (将先下载转为 Base64)
    image_url = "https://images.unsplash.com/photo-1575936123452-b67c3203c357" 
    print(f"正在获取图片并转换为 Base64: {image_url}")
    
    base64_image = encode_image_to_base64(image_url)
    
    if not base64_image:
        print("❌ 无法获取图片 Base64，跳过测试")
        return

    try:
        response = client.chat.completions.create(
            model=MODEL_NAME,
            messages=[
                {
                    "role": "user",
                    "content": [
                        {"type": "text", "text": "这张图片里有什么？请详细描述。"},
                        {
                            "type": "image_url",
                            "image_url": {
                                "url": f"data:image/jpeg;base64,{base64_image}"
                            }
                        },
                    ],
                }
            ],
        )
        print("\n[Kimi 回复]:")
        print(response.choices[0].message.content)
        print("\n[测试结果]: 图片理解成功 ✅")
    except Exception as e:
        print(f"\n[测试结果]: 图片理解失败 ❌\n错误信息: {e}")

def generate_multimodal_files():
    """
    生成多种格式的测试文件：txt, md, csv, pdf, docx, pptx, py, js
    """
    print("\n" + "="*50)
    print("正在生成测试文件...")
    print("="*50)

    files = []

    # 1. TXT
    txt_path = "test_gen.txt"
    with open(txt_path, "w", encoding="utf-8") as f:
        f.write("这是一个普通的文本文件。\nKimi 是 Moonshot AI 的助手。")
    files.append(txt_path)
    print(f"✅ 已生成: {txt_path}")

    # 2. Markdown
    md_path = "test_gen.md"
    with open(md_path, "w", encoding="utf-8") as f:
        f.write("# Markdown 测试\n\n- Kimi\n- Moonshot AI\n\n**加粗文本测试**")
    files.append(md_path)
    print(f"✅ 已生成: {md_path}")

    # 3. CSV
    csv_path = "test_gen.csv"
    with open(csv_path, "w", newline="", encoding="utf-8") as f:
        writer = csv.writer(f)
        writer.writerow(["Name", "Role", "Company"])
        writer.writerow(["Kimi", "Assistant", "Moonshot AI"])
    files.append(csv_path)
    print(f"✅ 已生成: {csv_path}")

    # 4. Code (Python)
    py_path = "test_gen.py"
    with open(py_path, "w", encoding="utf-8") as f:
        f.write("def hello():\n    print('Hello Kimi')\n\nif __name__ == '__main__':\n    hello()")
    files.append(py_path)
    print(f"✅ 已生成: {py_path}")

    # 5. Code (JavaScript)
    js_path = "test_gen.js"
    with open(js_path, "w", encoding="utf-8") as f:
        f.write("console.log('Hello Kimi');\nfunction greet(name) { return 'Hello ' + name; }")
    files.append(js_path)
    print(f"✅ 已生成: {js_path}")

    # 6. PDF (如果库存在)
    if FPDF:
        try:
            pdf_path = "test_gen.pdf"
            pdf = FPDF()
            pdf.add_page()
            pdf.set_font("Arial", size=12)
            pdf.cell(200, 10, txt="This is a PDF test file for Kimi.", ln=1, align="C")
            pdf.cell(200, 10, txt="Moonshot AI - File Understanding", ln=2, align="C")
            pdf.output(pdf_path)
            files.append(pdf_path)
            print(f"✅ 已生成: {pdf_path}")
        except Exception as e:
            print(f"❌ 生成 PDF 失败: {e}")
    else:
        print("⚠️  跳过 PDF 生成 (缺少 fpdf 库)")

    # 7. DOCX (如果库存在)
    if Document:
        try:
            docx_path = "test_gen.docx"
            doc = Document()
            doc.add_heading('Kimi DOCX Test', 0)
            doc.add_paragraph('这是一个由 Python 生成的 Word 文档。')
            doc.add_paragraph('用于测试文件上传与理解功能。')
            doc.save(docx_path)
            files.append(docx_path)
            print(f"✅ 已生成: {docx_path}")
        except Exception as e:
            print(f"❌ 生成 DOCX 失败: {e}")
    else:
        print("⚠️  跳过 DOCX 生成 (缺少 python-docx 库)")

    # 8. PPTX (如果库存在)
    if Presentation:
        try:
            # 8.1 PPTX
            pptx_path = "test_gen.pptx"
            prs = Presentation()
            slide_layout = prs.slide_layouts[0] # Title Slide
            slide = prs.slides.add_slide(slide_layout)
            title = slide.shapes.title
            subtitle = slide.placeholders[1]
            title.text = "Kimi PPTX Test"
            subtitle.text = "Generated by Python\nMoonshot AI"
            prs.save(pptx_path)
            files.append(pptx_path)
            print(f"✅ 已生成: {pptx_path}")

            # 8.2 PPT (注意: python-pptx 不支持保存为 .ppt 格式，这里仅作为模拟文件名)
            # 实际上 Kimi API 支持 ppt 和 pptx 格式上传。
            # 为了测试，我们简单复制 pptx 为 ppt 后缀 (仅作为文件扩展名测试，内容仍是 pptx 结构)
            # 在真实场景中，用户可能会上传旧版 Office 97-2003 的二进制 .ppt 文件
            ppt_path = "test_gen.ppt"
            with open(pptx_path, 'rb') as src, open(ppt_path, 'wb') as dst:
                dst.write(src.read())
            files.append(ppt_path)
            print(f"✅ 已生成: {ppt_path} (模拟)")

        except Exception as e:
            print(f"❌ 生成 PPTX/PPT 失败: {e}")
    else:
        print("⚠️  跳过 PPTX/PPT 生成 (缺少 python-pptx 库)")

    return files

def test_batch_file_understanding():
    """
    批量测试多种格式文件的上传与理解
    """
    print("\n" + "="*50)
    print("开始批量测试：多格式文件理解")
    print("="*50)

    # 1. 生成测试文件
    files = generate_multimodal_files()
    
    # 2. 逐个测试
    print("\n开始逐个上传并测试文件...")
    for file_path in files:
        test_file_understanding(file_path)
        # 避免请求过于频繁，稍微暂停
        time.sleep(1)

def test_file_understanding(file_path):
    """
    测试文件理解功能 (File Understanding)
    上传文件 -> 等待解析 -> 获取内容 -> 对话
    """
    print("\n" + "-"*30)
    print(f"正在测试文件: {file_path}")
    print("-" * 30)

    if not os.path.exists(file_path):
        print(f"文件不存在: {file_path}")
        return

    try:
        # 1. 上传文件
        print("1. 正在上传文件...")
        with open(file_path, "rb") as f:
            file_object = client.files.create(
                file=f,
                purpose="file-extract" # Moonshot 特有参数
            )
        file_id = file_object.id
        print(f"   文件上传成功，ID: {file_id}")

        # 2. 等待文件解析 (Moonshot 需要时间解析文件内容)
        # 注意：通常上传后需要一点时间才能获取内容，但在 SDK 中 content() 方法通常会处理
        # 或者我们可以尝试直接获取内容
        print("2. 获取文件解析内容...")
        # 简单重试机制
        content = ""
        for i in range(3):
            try:
                file_content = client.files.content(file_id=file_id)
                # client.files.content 返回的是 FileContent 对象，通常包含 text 属性
                # 注意：OpenAI SDK 标准返回可能是 bytes/text，Moonshot 返回的是解析后的文本
                # 这里假设直接返回文本内容
                content = file_content.text
                break
            except Exception as e:
                print(f"   等待解析中... ({e})")
                time.sleep(2)
        
        if not content:
            print("   无法获取文件内容，可能解析失败或超时。")
            return

        print(f"   成功获取内容 (前100字符): {content[:100]}...")

        # 3. 基于文件内容进行对话
        print("3. 发送给 Kimi 进行理解...")
        messages = [
            {
                "role": "system",
                "content": "你是 Kimi，请根据用户提供的文件内容回答问题。"
            },
            {
                "role": "system",
                "content": content # 将文件内容放入上下文
            },
            {
                "role": "user",
                "content": "请总结这个文件的主要内容。"
            }
        ]

        response = client.chat.completions.create(
            model=MODEL_NAME,
            messages=messages,
        )
        
        print("\n[Kimi 回复]:")
        print(response.choices[0].message.content)
        print("\n[测试结果]: 文件理解成功 ✅")

    except Exception as e:
        print(f"\n[测试结果]: 文件理解失败 ❌\n错误信息: {e}")

def test_video_understanding_real(video_file_path):
    """
    测试视频理解 (Video Understanding) - 使用 Kimi k2.5 上传视频
    步骤: 上传视频文件 -> 获取 file_id -> 构造 video_url (ms://<file_id>) -> 发送给 Kimi
    """
    print("\n" + "="*50)
    print("开始测试：视频理解 (Kimi k2.5 Video Upload)")
    print("="*50)

    if not os.path.exists(video_file_path):
        print(f"❌ 视频文件不存在: {video_file_path}")
        return

    try:
        # 1. 上传视频文件
        print(f"1. 正在上传视频文件: {video_file_path} ...")
        
        with open(video_file_path, "rb") as f:
            file_object = client.files.create(
                file=f,
                purpose="video" 
            )
        file_id = file_object.id
        print(f"   视频上传成功，ID: {file_id}")

        # 2. 构造请求
        # 格式: type="video_url", video_url={"url": "ms://file_id"}
        print("2. 发送给 Kimi 进行视频理解...")
        
        # 官方文档示例结构
        messages = [
            {
                "role": "user",
                "content": [
                    {"type": "video_url", "video_url": {"url": f"ms://{file_id}"}},
                    {"type": "text", "text": "请描述这个视频的内容。"}
                ]
            }
        ]

        # 必须使用支持 Vision 的模型，如 kimi-k2.5
        # 注意: message.content 字段类型为 List[Dict]
        response = client.chat.completions.create(
            model=MODEL_NAME, 
            messages=messages,
        )
        
        print("\n[Kimi 回复]:")
        print(response.choices[0].message.content)
        print("\n[测试结果]: 视频理解(上传方式)成功 ✅")

    except Exception as e:
        print(f"\n[测试结果]: 视频理解(上传方式)失败 ❌\n错误信息: {e}")


class AliyunASRCallback(RecognitionCallback):
    def __init__(self):
        self.transcribed_text = ""

    def on_open(self) -> None:
        print("   [Aliyun ASR] 连接建立")

    def on_close(self) -> None:
        print("   [Aliyun ASR] 连接关闭")

    def on_complete(self) -> None:
        print("   [Aliyun ASR] 识别完成")

    def on_error(self, result: RecognitionResult) -> None:
        print(f"   [Aliyun ASR] 错误: {result}")

    def on_event(self, result: RecognitionResult) -> None:
        sentence = result.get_sentence()
        if 'text' in sentence:
            text = sentence['text']
            if result.is_sentence_end(sentence):
                print(f"   [Aliyun ASR] 识别结果: {text}")
                self.transcribed_text += text

def test_audio_understanding_aliyun(audio_file_path):
    """
    测试音频理解 (Audio Understanding) - 使用阿里云 DashScope FunASR
    步骤: 音频文件 -> Aliyun FunASR -> 文本 -> Kimi
    """
    print("\n" + "="*50)
    print("开始测试：音频理解 (Aliyun FunASR + Kimi)")
    print("="*50)
    
    # 检查 DashScope 库
    if not dashscope:
        print("❌ 未安装 dashscope 库，跳过测试。请运行: pip install dashscope")
        return

    # 检查 API Key
    if not DASHSCOPE_API_KEY:
        print("⚠️  未检测到 DASHSCOPE_API_KEY 环境变量。")
        print("   请设置 DASHSCOPE_API_KEY 以使用阿里云语音识别服务。")
        print("   本次测试将跳过实际 ASR 调用，仅做提示。")
        return
    
    dashscope.api_key = DASHSCOPE_API_KEY
    
    # 检查音频文件
    if not os.path.exists(audio_file_path):
        print(f"❌ 音频文件不存在: {audio_file_path}")
        # 尝试生成一个
        try:
            # 简单的生成逻辑，或者直接提示用户
            print("   (提示: 请确保目录下存在有效的音频文件，例如 test_audio.wav)")
        except:
            pass
        return

    print(f"正在处理音频文件: {audio_file_path}")
    print("1. 调用阿里云 FunASR 进行语音转文字...")

    callback = AliyunASRCallback()
    
    # 使用 fun-asr-realtime-v1 模型 (或者 paraformer-realtime-v1)
    # 注意: 阿里云推荐实时语音使用 paraformer-realtime-v1
    # 官方文档中提到的 fun-asr-realtime 也可用
    recognition = Recognition(
        model='paraformer-realtime-v1', # 推荐模型
        format='pcm', # 假设是 PCM 或者 wav (SDK会自动处理部分格式，但最好是 PCM/WAV)
        # 如果是 wav 文件，SDK 通常可以处理，但最好指定 sample_rate
        sample_rate=16000, 
        callback=callback
    )

    recognition.start()

    # 读取并发送音频数据
    # 注意: 这里的实现是模拟实时流，将文件分块发送
    try:
        # 尝试使用 soundfile 或 wave 读取
        # 这里为了简单，假设是 wav/pcm 16k mono
        # 如果是 wav，需要跳过 header 或者直接读取 data
        import wave
        with wave.open(audio_file_path, 'rb') as wf:
            # 检查参数
            channels = wf.getnchannels()
            rate = wf.getframerate()
            if rate != 16000:
                print(f"⚠️  警告: 音频采样率是 {rate}，模型通常需要 16000。可能会识别失败。")
            
            chunk_size = 3200 # 100ms for 16k
            # 简单重采样逻辑：如果采样率是 44100，每隔 n 个点取一个，或者使用 audioop
            # 为了更好的效果，建议使用 audioop 进行降采样
            
            data = wf.readframes(chunk_size)
            while len(data) > 0:
                # 如果采样率不匹配，尝试转换 (仅支持简单的降采样，例如 48k->16k, 32k->16k)
                # 44.1k -> 16k 比较复杂，这里简单处理，如果需要高质量建议用 ffmpeg 预处理
                if rate != 16000:
                    try:
                        # 简单的降采样：每隔 step 取一个点
                        step = rate // 16000
                        if step > 1:
                            # 假设是 16bit PCM (2 bytes per sample)
                            import struct
                            num_samples = len(data) // 2
                            if num_samples > 0:
                                samples = struct.unpack(f"{num_samples}h", data)
                                resampled_samples = samples[::step]
                                data = struct.pack(f"{len(resampled_samples)}h", *resampled_samples)
                    except Exception:
                        pass # 转换失败则发送原始数据

                recognition.send_audio_frame(data)
                # 模拟实时发送的间隔 (可选)
                time.sleep(0.01) 
                data = wf.readframes(chunk_size)
        
        recognition.stop()
    except Exception as e:
        print(f"❌ 发送音频数据失败: {e}")
        recognition.stop()
        return

    transcribed_text = callback.transcribed_text
    print(f"\n   ASR 转录结果: {transcribed_text}")
    
    if not transcribed_text:
        print("⚠️  未能识别出文本 (可能是音频为空或格式不支持)。")
        # 为了流程继续，使用默认文本 (如果用户同意)
        # transcribed_text = "这是测试文本"
        return

    print("2. 发送转录文本给 Kimi...")
    try:
        response = client.chat.completions.create(
            model=MODEL_NAME,
            messages=[
                {
                    "role": "system",
                    "content": "你是一个语音助手，请根据用户的语音转录内容进行回复。"
                },
                {
                    "role": "user",
                    "content": transcribed_text
                }
            ],
        )
        print("\n[Kimi 回复]:")
        print(response.choices[0].message.content)
        print("\n[测试结果]: 音频理解(Aliyun ASR)成功 ✅")
    except Exception as e:
        print(f"\n[测试结果]: 音频理解(Aliyun ASR)失败 ❌\n错误信息: {e}")

def test_audio_understanding():
    """
    测试音频理解 (Audio Understanding) - 使用阿里云 DashScope FunASR
    步骤: 音频文件 -> Aliyun FunASR -> 文本 -> Kimi
    """
    print("\n" + "="*50)
    print("开始测试：音频理解 (Aliyun FunASR + Kimi)")
    print("="*50)
    
    # 检查 DashScope 库
    if not dashscope:
        print("❌ 未安装 dashscope 库，跳过测试。请运行: pip install dashscope")
        return

    # 检查 API Key
    if not DASHSCOPE_API_KEY:
        print("⚠️  未检测到 DASHSCOPE_API_KEY 环境变量。")
        print("   请设置 DASHSCOPE_API_KEY 以使用阿里云语音识别服务。")
        print("   本次测试将跳过实际 ASR 调用，仅做提示。")
        return
    
    dashscope.api_key = DASHSCOPE_API_KEY
    
    # 检查音频文件
    if not os.path.exists(audio_file_path):
        print(f"❌ 音频文件不存在: {audio_file_path}")
        # 尝试生成一个
        try:
            # 简单的生成逻辑，或者直接提示用户
            print("   (提示: 请确保目录下存在有效的音频文件，例如 test_audio.wav)")
        except:
            pass
        return

    print(f"正在处理音频文件: {audio_file_path}")
    print("1. 调用阿里云 FunASR 进行语音转文字...")

    callback = AliyunASRCallback()
    
    # 使用 fun-asr-realtime-v1 模型 (或者 paraformer-realtime-v1)
    # 注意: 阿里云推荐实时语音使用 paraformer-realtime-v1
    # 官方文档中提到的 fun-asr-realtime 也可用
    recognition = Recognition(
        model='paraformer-realtime-v1', # 推荐模型
        format='pcm', # 假设是 PCM 或者 wav (SDK会自动处理部分格式，但最好是 PCM/WAV)
        # 如果是 wav 文件，SDK 通常可以处理，但最好指定 sample_rate
        sample_rate=16000, 
        callback=callback
    )

    recognition.start()

    # 读取并发送音频数据
    # 注意: 这里的实现是模拟实时流，将文件分块发送
    try:
        # 尝试使用 soundfile 或 wave 读取
        # 这里为了简单，假设是 wav/pcm 16k mono
        # 如果是 wav，需要跳过 header 或者直接读取 data
        import wave
        with wave.open(audio_file_path, 'rb') as wf:
            # 检查参数
            channels = wf.getnchannels()
            rate = wf.getframerate()
            if rate != 16000:
                print(f"⚠️  警告: 音频采样率是 {rate}，模型通常需要 16000。可能会识别失败。")
            
            chunk_size = 3200 # 100ms for 16k
            data = wf.readframes(chunk_size)
            while len(data) > 0:
                recognition.send_audio_frame(data)
                # 模拟实时发送的间隔 (可选)
                time.sleep(0.01) 
                data = wf.readframes(chunk_size)
        
        recognition.stop()
    except Exception as e:
        print(f"❌ 发送音频数据失败: {e}")
        recognition.stop()
        return

    transcribed_text = callback.transcribed_text
    print(f"\n   ASR 转录结果: {transcribed_text}")
    
    if not transcribed_text:
        print("⚠️  未能识别出文本 (可能是音频为空或格式不支持)。")
        # 为了流程继续，使用默认文本 (如果用户同意)
        # transcribed_text = "这是测试文本"
        return

    print("2. 发送转录文本给 Kimi...")
    try:
        response = client.chat.completions.create(
            model=MODEL_NAME,
            messages=[
                {
                    "role": "system",
                    "content": "你是一个语音助手，请根据用户的语音转录内容进行回复。"
                },
                {
                    "role": "user",
                    "content": transcribed_text
                }
            ],
        )
        print("\n[Kimi 回复]:")
        print(response.choices[0].message.content)
        print("\n[测试结果]: 音频理解(Aliyun ASR)成功 ✅")
    except Exception as e:
        print(f"\n[测试结果]: 音频理解(Aliyun ASR)失败 ❌\n错误信息: {e}")

def main():
    print("Moonshot AI (Kimi) 多模态 API 测试脚本")
    print(f"当前使用的模型: {MODEL_NAME}")
    print("注意: 请确保已设置 MOONSHOT_MOONSHOT_API_KEY 环境变量")
    
    # 1. 测试图片
    test_image_understanding()
    
    # 2. 测试文件 (多格式批量测试)
    # 创建并测试多种格式文件 (TXT, MD, PDF, DOCX, PPTX, CSV, CODE)
    test_batch_file_understanding()
    
    # 3. 测试视频
    local_video_path = "test_video.mp4"
    if os.path.exists(local_video_path):
        test_video_understanding_real(local_video_path)
    else:
        print(f"⚠️  未检测到本地视频文件: {local_video_path}")
        print("   跳过真实上传测试。")
    
    # 4. 测试音频 (阿里云ASR)
    # 直接使用 Aliyun ASR 进行测试
    # 需要确保存在 test_audio.wav 且格式正确 (16k, mono, 16bit)
    if not os.path.exists("test_audio.wav"):
        print("❌ 未找到 test_audio.wav，正在尝试生成...")
        try:
            # 尝试调用 generate_audio.py 生成
            import generate_audio
            generate_audio.create_sine_wave("test_audio.wav", framerate=16000)
            print("✅ 已生成 test_audio.wav")
        except ImportError:
            print("❌ 无法生成音频文件，跳过音频测试。")
    
    if os.path.exists("test_audio.wav"):
        test_audio_understanding_aliyun("test_audio.wav")
    else:
        print("❌ 音频文件 test_audio.wav 不存在，跳过阿里云 ASR 测试。")

if __name__ == "__main__":
    # 如果通过命令行参数传入 --no-run，则不运行测试
    if len(sys.argv) > 1 and sys.argv[1] == "--no-run":
        print("已生成脚本，但根据指令跳过运行测试。")
        sys.exit(0)
    main()

